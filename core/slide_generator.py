import os
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from core.models import Slide
from core.utils import setup_logger

logger = setup_logger(__name__)

# Configurações de cores e tamanhos (Paleta do Usuário)
PRIMARY = RGBColor(0x00, 0x8A, 0x7A)       # #008A7A
PRIMARY_DARK = RGBColor(0x00, 0x4D, 0x40)  # Verde ainda mais escuro
ACCENT = RGBColor(0xF7, 0x94, 0x1D)        # #F7941D
BG = RGBColor(0xFA, 0xFA, 0xFA)            # #FAFAFA (quase branco)
TEXT = RGBColor(0x1F, 0x1F, 0x1F)          # #1F1F1F
DARK_GREY = RGBColor(0x33, 0x33, 0x33)

FONT_NAME = "Roboto Slab"

class SlideGenerator:
    """Classe responsável por transformar uma lista de Slides num arquivo .pptx"""
    
    def __init__(self, slides: List[Slide]):
        self.slides = slides
        self.prs = Presentation()
        # Widescreen 16:9
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.blank_layout = self.prs.slide_layouts[6]
        
    def _create_base_slide(self, slide_model: Slide):
        """Cria o slide e configura o background e o logo."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # Fundo geral (Background clean)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG
        
        # Inserir Logo no canto inferior direito
        base_dir = os.path.dirname(os.path.dirname(__file__))
        logo_path = os.path.join(base_dir, "fonte", "logo.png")
        if os.path.exists(logo_path):
            logo_height = Inches(1.2)
            # Como não sabemos a largura, colocamos uma posição aproximada ancorada pela direita. 
            # O add_picture ajusta a proporção automaticamente se passarmos apenas height
            logo_left = self.prs.slide_width - Inches(3.0) 
            logo_top = Inches(0.3)
            try:
                slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)
            except Exception as e:
                logger.warning(f"Não foi possível inserir o logo no slide: {e}")
                
        return slide

    def _add_content_text(self, slide, slide_model: Slide):
        """Adiciona a caixa de texto de conteúdo de acordo com o tipo de slide."""
        
        if slide_model.is_cover or slide_model.is_section:
            # Título principal centralizado
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), self.prs.slide_width - Inches(2), Inches(1.5))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p = tf_title.paragraphs[0]
            p.text = slide_model.title
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.name = FONT_NAME
            p.font.color.rgb = PRIMARY_DARK
            
            # Linha decorativa
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(5.66), Inches(3.9), Inches(2), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT
            line.line.fill.background()
            
            # Subtítulo ou conteúdo centralizado
            if slide_model.content:
                content_box = slide.shapes.add_textbox(Inches(1), Inches(4.1), self.prs.slide_width - Inches(2), Inches(1.5))
                tf_content = content_box.text_frame
                tf_content.word_wrap = True
                p2 = tf_content.paragraphs[0]
                p2.text = slide_model.content
                p2.alignment = PP_ALIGN.CENTER
                p2.font.size = Pt(32)
                p2.font.name = FONT_NAME
                p2.font.color.rgb = TEXT
                
            # Data no rodapé
            if slide_model.date:
                date_box = slide.shapes.add_textbox(Inches(2), Inches(6.5), self.prs.slide_width - Inches(4), Inches(0.5))
                tf_date = date_box.text_frame
                p3 = tf_date.paragraphs[0]
                p3.text = slide_model.date
                p3.alignment = PP_ALIGN.CENTER
                p3.font.size = Pt(20)
                p3.font.bold = True
                p3.font.name = FONT_NAME
                p3.font.color.rgb = DARK_GREY
                
        else:
            # Layout Padrão Clean
            
            # Título no topo esquerdo
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), self.prs.slide_width - Inches(3.5), Inches(1.2))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_model.title
            p_title.alignment = PP_ALIGN.LEFT
            p_title.font.size = Pt(36)
            p_title.font.bold = True
            p_title.font.name = FONT_NAME
            p_title.font.color.rgb = PRIMARY_DARK
            
            # Linha laranja fina abaixo do título
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), self.prs.slide_width - Inches(1.6), Inches(0.04)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT
            line.line.fill.background()
            
            # Área de Texto Principal
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), self.prs.slide_width - Inches(1.6), self.prs.slide_height - Inches(2.2))
            text_frame = txBox.text_frame
            text_frame.word_wrap = True
            
            paragraphs_text = slide_model.content.split('\n')
            
            for i, pt_text in enumerate(paragraphs_text):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                    
                if slide_model.is_bullet and pt_text.strip().startswith('•'):
                    p.text = pt_text.replace('•', '').strip()
                    p.level = 0
                else:
                    p.text = pt_text
                    
                p.font.color.rgb = TEXT
                p.font.name = FONT_NAME
                p.line_spacing = 1.3
                p.font.size = Pt(28)
                p.alignment = PP_ALIGN.JUSTIFY
                
        return slide

    def generate(self, output_path: str):
        logger.info("Iniciando geração de slides PPTX...")
        
        for s in self.slides:
            slide = self._create_base_slide(s)
            self._add_content_text(slide, s)

        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        logger.info(f"Slides salvos com sucesso em: {output_path}")
        return output_path
